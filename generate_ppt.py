from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()

    # Définition des slides (Titre, Contenu, Notes orales)
    slides_data = [
        {
            "title": "ESME : Moteur de Recommandation Musicale Hybride",
            "content": [
                "Architecture, Algorithmique Vectorielle et Intégration API Spotify",
                "",
                "Présenté par : [Ton Nom]",
                "École : [Ton École]",
                "Date : Février 2026"
            ],
            "notes": "Bonjour à tous. Je vous présente aujourd'hui le projet ESME. Au-delà d'une simple interface web, ce projet visait à concevoir et implémenter un moteur de recommandation hybride, capable de résoudre le problème de la découverte musicale par l'analyse de signaux audio."
        },
        {
            "title": "Problématique Métier : Le Paradoxe du Choix",
            "content": [
                "Contexte : Catalogues de streaming > 100 millions de titres.",
                "Problème 1 : Surcharge Informationnelle (Utilisateur perdu).",
                "Problème 2 : Le 'Cold Start' (Comment recommander sans historique ?).",
                "Objectif : Maximiser la Pertinence (Accuracy) et la Diversité (Serendipity)."
            ],
            "notes": "La problématique est double. 1: La surcharge informationnelle. 2: Le problème technique du 'Cold Start'. Mon système vise à optimiser deux métriques souvent opposées : la pertinence mathématique et la sérendipité."
        },
        {
            "title": "La Solution ESME : Approche Hybride",
            "content": [
                "Application Web Interactive (Frontend Jinja2/JS).",
                "3 Stratégies de Recommandation :",
                "   1. Simple (Déduplication).",
                "   2. Content-Based (Analyse Audio).",
                "   3. Hybride (Fusion Pondérée).",
                "Intégration complète à l'écosystème Spotify."
            ],
            "notes": "J'ai développé une solution permettant 3 stratégies. L'innovation principale réside dans l'approche hybride, fusionnant l'analyse du signal audio et les données sociales."
        },
        {
            "title": "Architecture Technique Globale",
            "content": [
                "Backend : Python / FastAPI (Asynchrone).",
                "Design Pattern : MVC (Modèle-Vue-Contrôleur).",
                "Data Processing : Pandas (DataFrames) & Numpy (Calcul matriciel).",
                "Protocole : Communication RESTful avec Spotify Web API."
            ],
            "notes": "Architecture Python moderne. FastAPI pour l'asynchronisme (crucial pour les API externes). Pattern MVC pour séparer la logique métier de la présentation."
        },
        {
            "title": "Ingénierie des Données (Feature Engineering)",
            "content": [
                "Source : Extraction des 'Audio Features' via API.",
                "Espace Vectoriel à N-Dimensions :",
                "   - Danceability, Energy, Valence (0.0 - 1.0)",
                "   - Acousticness, Instrumentalness...",
                "Normalisation : Données bornées pour éviter les biais d'échelle."
            ],
            "notes": "Le cœur du système. Chaque morceau est un vecteur dans un espace N-dimensionnel. J'utilise des descripteurs de haut niveau comme la Valence ou l'Énergie pour calculer les similarités."
        },
        {
            "title": "Modélisation Mathématique : Cosine Similarity",
            "content": [
                "Métrique choisie : Similarité Cosinus.",
                "Formule : Cos(theta) = (A . B) / (||A|| . ||B||)",
                "Avantages :",
                "   - Mesure l'orientation (l'angle) et non la magnitude.",
                "   - Efficace en haute dimensionnalité.",
                "   - Résultat normalisé entre 0 (différent) et 1 (identique)."
            ],
            "notes": "J'utilise la Similarité Cosinus. Elle mesure l'angle entre les vecteurs, donnant un score de similarité pure indépendant de la 'popularité' (magnitude). C'est le moteur du filtrage par contenu."
        },
        {
            "title": "Algorithme 1 : Filtrage Collaboratif (Social)",
            "content": [
                "Principe : 'User-Based Collaborative Filtering'.",
                "Méthode : Exploitation des clusters utilisateurs Spotify.",
                "Données : Matrices d'interactions (Top Tracks).",
                "Limite : Biais de popularité (Tendance à recommander du Mainstream)."
            ],
            "notes": "Le filtrage collaboratif ne regarde pas le signal audio, mais les comportements. On utilise l'API Spotify comme proxy de données massives pour trouver ce que les profils similaires aiment."
        },
        {
            "title": "Algorithme 2 : Content-Based Filtering (Signal)",
            "content": [
                "Principe : Recommandation basée sur les propriétés intrinsèques.",
                "Logique Vectorielle :",
                "   1. Calcul du Centroid (Vecteur Moyen) de l'utilisateur.",
                "   2. Calcul des distances Cosinus avec les candidats.",
                "   3. Ranking par score décroissant.",
                "Avantage : Résout le problème du Cold Start."
            ],
            "notes": "Le Content-Based, implémenté manuellement. Je calcule le vecteur moyen des goûts de l'utilisateur, puis je scanne la base pour trouver les vecteurs géométriquement proches."
        },
        {
            "title": "L'Algorithme Hybride ESME : Le Re-Ranking",
            "content": [
                "Stratégie : Weighted Hybrid / Switching.",
                "Workflow :",
                "   1. Génération de Candidats (via Collaboratif).",
                "   2. Enrichissement (Extraction Audio Features).",
                "   3. Filtrage Vectoriel (Cosine Similarity vs User Seeds).",
                "   4. Re-Ranking (Tri final par cohérence audio)."
            ],
            "notes": "Mon algo final est hybride. Il utilise le collaboratif pour le 'Sourcing' large, et le Content-Based pour le 'Filtering' fin. Ce Re-ranking assure une playlist populaire ET cohérente."
        },
        {
            "title": "Implémentation Python & Performance",
            "content": [
                "Stack : Pandas (DataFrames) + Scikit-Learn.",
                "Optimisation :",
                "   - Vectorisation des opérations (Numpy).",
                "   - Complexité O(N*M) optimisée vs Boucles natives.",
                "   - Utilisation de sklearn.metrics.pairwise_distances."
            ],
            "notes": "Implémentation via Pandas/Scikit-Learn. Point critique : optimisation des calculs via vectorisation Numpy pour éviter les boucles Python lentes et assurer le temps réel."
        },
        {
            "title": "Pipeline ETL Temps Réel",
            "content": [
                "Extract : Auth OAuth2 + Fetch Top Tracks + Audio Features.",
                "Transform : Normalisation, Vectorisation, Moyennes.",
                "Load : Création Playlist, Ajout items, Upload Cover.",
                "Robustesse : Gestion des Rate Limits et Exceptions API."
            ],
            "notes": "C'est un pipeline ETL temps réel. Extraction des métadonnées, transformation en vecteurs, application du modèle, et chargement sur Spotify. Tout ça en quelques secondes."
        },
        {
            "title": "Sécurité : OAuth 2.0 Authorization Code",
            "content": [
                "Protocole : OAuth 2.0 Standard.",
                "Flux : Authorization Code Flow (Sécurisé serveur).",
                "Scopes : 'user-top-read', 'playlist-modify-public'.",
                "Sécurité : Pas de stockage de mot de passe (Tokens éphémères)."
            ],
            "notes": "Sécurité via OAuth 2.0. Aucun mot de passe stocké. Utilisation de Tokens d'accès éphémères et de Refresh Tokens. Scopes limités au strict nécessaire."
        },
        {
            "title": "Limites et Perspectives (Master)",
            "content": [
                "Limites :",
                "   - Dépendance API Spotify.",
                "   - Modèle statique (pas de notion de séquence).",
                "Perspectives :",
                "   - Deep Learning (RNN/LSTM) pour l'ordre séquentiel.",
                "   - NLP sur les paroles pour l'analyse sémantique."
            ],
            "notes": "Limites : dépendance API. Avenir : Intégrer des Réseaux de Neurones Récurrents (RNN) pour gérer la SÉQUENCE des titres (le mix) et non plus juste la liste."
        },
        {
            "title": "Conclusion",
            "content": [
                "Projet validé de bout en bout (Full Stack).",
                "Application concrète de l'Algèbre Linéaire (Vecteurs).",
                "Architecture Web robuste et Sécurisée.",
                "Démonstration fonctionnelle (Playlist 'esme').",
                "Merci de votre attention."
            ],
            "notes": "Le projet prouve qu'on peut intégrer des maths (espaces vectoriels) dans une app web moderne. C'est fonctionnel, sécurisé et scalable. Merci."
        }
    ]

    # Création des slides
    for slide_info in slides_data:
        # Layout 1 est souvent "Titre et Contenu"
        slide_layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(slide_layout)
        
        # Titre
        title = slide.shapes.title
        title.text = slide_info["title"]
        
        # Contenu (Bullet points)
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.text = slide_info["content"][0] # Première ligne
        
        for line in slide_info["content"][1:]:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0 # Niveau d'indentation principal
            
        # Ajout des notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = slide_info["notes"]

    # Sauvegarde
    prs.save('Presentation_ESME_Master.pptx')
    print("Présentation générée avec succès : Presentation_ESME_Master.pptx")

if __name__ == "__main__":
    create_presentation()
    